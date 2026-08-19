import asyncio
import copy
import os
import shutil
import tempfile
from pathlib import Path
from typing import Any
from urllib.parse import quote

import fitz
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.shared import Inches, Pt
from fastapi import FastAPI, Header, HTTPException
from openpyxl import Workbook, load_workbook
from openpyxl.chart import BarChart, Reference
from openpyxl.styles import Alignment, Font, PatternFill
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt
from pydantic import BaseModel
from pypdf import PdfReader, PdfWriter

ROOT = Path(os.getenv("WORKSPACE_ROOT", "/workspace")).resolve()
TEMPLATES = Path(os.getenv("TEMPLATE_ROOT", "/templates")).resolve()
SHARED_SECRET = os.getenv("AGENT_SHARED_SECRET", "")
app = FastAPI(title="Offline AI Office Worker", version="1.0.0")


class ToolRequest(BaseModel):
    name: str
    arguments: dict[str, Any] = {}


def authorize(value: str | None) -> None:
    if SHARED_SECRET and value != f"Bearer {SHARED_SECRET}":
        raise HTTPException(status_code=401, detail="invalid API key")


def safe_path(value: str, must_exist: bool = False) -> Path:
    path = (ROOT / value.lstrip("/\\")).resolve(strict=must_exist)
    if path != ROOT and ROOT not in path.parents:
        raise ValueError("path escapes /workspace")
    return path


def safe_template(kind: str, value: str) -> Path:
    base = (TEMPLATES / kind).resolve()
    path = (base / value).resolve(strict=True)
    if path != base and base not in path.parents:
        raise ValueError("template path escapes template directory")
    return path


def output_path(value: str, suffix: str) -> Path:
    raw = Path(value)
    name = raw.name if raw.suffix.lower() == suffix else raw.name + suffix
    relative = raw.parent / name if str(raw.parent) not in {".", ""} else Path("outputs") / name
    path = safe_path(str(relative))
    path.parent.mkdir(parents=True, exist_ok=True)
    return path


def artifact(path: Path) -> dict[str, Any]:
    relative = path.relative_to(ROOT).as_posix()
    return {"path": relative, "url": f"/artifacts/{quote(relative)}", "size": path.stat().st_size}


def unique_sheet_name(book: Workbook, requested: str) -> str:
    base = requested[:31] or "Sheet"
    candidate = base
    index = 2
    while candidate in book.sheetnames:
        suffix = f"_{index}"
        candidate = f"{base[:31 - len(suffix)]}{suffix}"
        index += 1
    return candidate


async def libreoffice_convert(path: Path, target: str = "pdf") -> Path:
    out_dir = path.parent
    profile = Path(tempfile.mkdtemp(prefix="lo-profile-"))
    proc = await asyncio.create_subprocess_exec("libreoffice", f"-env:UserInstallation=file://{profile}", "--headless", "--nologo", "--nodefault", "--nolockcheck", "--convert-to", target, "--outdir", str(out_dir), str(path), stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE)
    try:
        stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=180)
    finally:
        shutil.rmtree(profile, ignore_errors=True)
    converted = out_dir / f"{path.stem}.{target.split(':')[0]}"
    if proc.returncode != 0 or not converted.exists() or converted.stat().st_size == 0:
        raise RuntimeError(f"LibreOffice conversion failed: {stdout.decode(errors='replace')} {stderr.decode(errors='replace')}")
    return converted


async def validate_office(path: Path, kind: str) -> dict[str, Any]:
    pdf = await libreoffice_convert(path)
    with fitz.open(pdf) as opened:
        pages = opened.page_count
    return {
        "validation": {kind: "PASS", "libreoffice_pdf": "PASS", "pdf_pages": pages},
        "pdf": artifact(pdf),
    }


def chinese_font(run, name: str = "Noto Sans CJK SC", size: int = 11, bold: bool = False):
    run.font.name = name
    run._element.rPr.rFonts.set(qn("w:eastAsia"), name)
    run.font.size = Pt(size)
    run.bold = bold


async def create_docx(a: dict[str, Any]) -> dict[str, Any]:
    path = output_path(a["path"], ".docx")
    template = a.get("template")
    doc = Document(str(safe_template("word", template))) if template else Document()
    title = doc.add_heading(level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    chinese_font(title.add_run(a.get("title", "文档")), size=20, bold=True)
    for text in a.get("paragraphs", []):
        paragraph = doc.add_paragraph()
        chinese_font(paragraph.add_run(str(text)), size=11)
    for item in a.get("lists", []):
        paragraph = doc.add_paragraph(style="List Number" if item.get("ordered") else "List Bullet")
        chinese_font(paragraph.add_run(str(item.get("text", ""))), size=int(item.get("size", 11)))
    table_data = a.get("table") or []
    if table_data:
        cols = max(len(row) for row in table_data)
        table = doc.add_table(rows=len(table_data), cols=cols)
        table.style = "Table Grid"
        for row_index, row in enumerate(table_data):
            for col_index, value in enumerate(row):
                cell = table.cell(row_index, col_index)
                cell.text = ""
                chinese_font(cell.paragraphs[0].add_run(str(value)), bold=row_index == 0)
    for item in a.get("images", []):
        image = safe_path(item["path"], True)
        doc.add_picture(str(image), width=Inches(float(item.get("width_inches", 6))))
    section = doc.sections[0]
    if a.get("header"):
        chinese_font(section.header.paragraphs[0].add_run(a["header"]), size=9)
    if a.get("footer"):
        chinese_font(section.footer.paragraphs[0].add_run(a["footer"]), size=9)
    if a.get("page_break"):
        doc.add_page_break()
    doc.save(path)
    Document(path)
    checked = await validate_office(path, "python_docx")
    return {"ok": True, "artifact": artifact(path), **checked}


def replace_docx_text(doc: Document, replacements: dict[str, Any]) -> None:
    def replace_paragraph(paragraph) -> None:
        if not replacements or not paragraph.text:
            return
        original = paragraph.text
        updated = original
        for old, new in replacements.items():
            updated = updated.replace(str(old), str(new))
        if updated != original:
            for run in paragraph.runs:
                run.text = ""
            run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
            run.text = updated
            chinese_font(run)

    for paragraph in doc.paragraphs:
        replace_paragraph(paragraph)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    replace_paragraph(paragraph)


async def edit_docx(a: dict[str, Any]) -> dict[str, Any]:
    source = safe_path(a.get("source") or a["path"], True)
    target = output_path(a.get("output") or a.get("path") or source.name, ".docx")
    doc = Document(source)
    replace_docx_text(doc, a.get("replace", {}))
    for text in a.get("append_paragraphs", []):
        chinese_font(doc.add_paragraph().add_run(str(text)))
    for item in a.get("append_lists", []):
        paragraph = doc.add_paragraph(style="List Number" if item.get("ordered") else "List Bullet")
        chinese_font(paragraph.add_run(str(item.get("text", ""))))
    table_data = a.get("append_table") or []
    if table_data:
        cols = max(len(row) for row in table_data)
        table = doc.add_table(rows=len(table_data), cols=cols)
        table.style = "Table Grid"
        for row_index, row in enumerate(table_data):
            for col_index, value in enumerate(row):
                table.cell(row_index, col_index).text = str(value)
    if a.get("page_break"):
        doc.add_page_break()
    doc.save(target)
    Document(target)
    checked = await validate_office(target, "python_docx")
    return {"ok": True, "artifact": artifact(target), **checked}


async def read_docx(a: dict[str, Any]) -> dict[str, Any]:
    path = safe_path(a["path"], True)
    doc = Document(path)
    tables = [[[cell.text for cell in row.cells] for row in table.rows] for table in doc.tables]
    return {"ok": True, "paragraphs": [p.text for p in doc.paragraphs], "tables": tables}


def normalize_rows(value: Any) -> list[list[Any]]:
    if isinstance(value, list):
        if not value:
            return []
        if all(isinstance(row, dict) for row in value):
            headers = list(value[0])
            return [headers] + [[row.get(header) for header in headers] for row in value]
        return [row if isinstance(row, list) else [row] for row in value]
    raise ValueError("sheet data must be an array")


async def create_xlsx(a: dict[str, Any]) -> dict[str, Any]:
    path = output_path(a["path"], ".xlsx")
    template = a.get("template")
    book = load_workbook(safe_template("excel", template)) if template else Workbook()
    if not template:
        book.remove(book.active)
    for sheet_name, raw_rows in a["sheets"].items():
        requested = str(sheet_name)[:31]
        sheet = book[requested] if requested in book.sheetnames else book.create_sheet(requested)
        rows = normalize_rows(raw_rows)
        for row in rows:
            sheet.append(row)
        sheet.freeze_panes = "A2"
        sheet.auto_filter.ref = sheet.dimensions
        for cell in sheet[1]:
            cell.font = Font(name="Noto Sans CJK SC", bold=True, color="FFFFFF")
            cell.fill = PatternFill("solid", fgColor="2F5597")
            cell.alignment = Alignment(horizontal="center")
        for column in range(1, sheet.max_column + 1):
            max_len = max((len(str(sheet.cell(row, column).value or "")) for row in range(1, sheet.max_row + 1)), default=8)
            sheet.column_dimensions[get_column_letter(column)].width = min(max(max_len * 1.5, 10), 60)
        for chart_spec in a.get("charts", []):
            if chart_spec.get("sheet") != sheet_name or sheet.max_row < 2 or sheet.max_column < 2:
                continue
            chart = BarChart()
            chart.title = chart_spec.get("title", "汇总图表")
            data = Reference(sheet, min_col=int(chart_spec.get("value_column", 2)), min_row=1, max_row=sheet.max_row)
            cats = Reference(sheet, min_col=int(chart_spec.get("category_column", 1)), min_row=2, max_row=sheet.max_row)
            chart.add_data(data, titles_from_data=True)
            chart.set_categories(cats)
            sheet.add_chart(chart, chart_spec.get("anchor", "E2"))
    for spec in a.get("merged_cells", []):
        book[str(spec["sheet"])[:31]].merge_cells(str(spec["range"]))
    book.save(path)
    checked = load_workbook(path, data_only=False)
    sheet_names = checked.sheetnames
    checked.close()
    rendered = await validate_office(path, "openpyxl_reload")
    rendered["validation"]["sheets"] = sheet_names
    return {"ok": True, "artifact": artifact(path), **rendered}


async def edit_xlsx(a: dict[str, Any]) -> dict[str, Any]:
    source = safe_path(a.get("source") or a["path"], True)
    target = output_path(a.get("output") or a.get("path") or source.name, ".xlsx")
    book = load_workbook(source)
    for name, rows in a.get("create_sheets", {}).items():
        sheet = book.create_sheet(unique_sheet_name(book, str(name)))
        for row in normalize_rows(rows):
            sheet.append(row)
    for name, rows in a.get("append_rows", {}).items():
        sheet = book[str(name)]
        for row in normalize_rows(rows):
            sheet.append(row)
    for item in a.get("set_cells", []):
        book[str(item["sheet"])][str(item["cell"])] = item.get("value")
    for item in a.get("formulas", []):
        formula = str(item["formula"])
        book[str(item["sheet"])][str(item["cell"])] = formula if formula.startswith("=") else f"={formula}"
    for spec in a.get("merged_cells", []):
        book[str(spec["sheet"])].merge_cells(str(spec["range"]))
    for name in a.get("delete_sheets", []):
        if str(name) in book.sheetnames and len(book.sheetnames) > 1:
            book.remove(book[str(name)])
    book.save(target)
    checked = load_workbook(target, data_only=False)
    sheet_names = checked.sheetnames
    checked.close()
    rendered = await validate_office(target, "openpyxl_reload")
    rendered["validation"]["sheets"] = sheet_names
    return {"ok": True, "artifact": artifact(target), **rendered}


async def merge_xlsx(a: dict[str, Any]) -> dict[str, Any]:
    target = output_path(a["path"], ".xlsx")
    output = Workbook()
    output.remove(output.active)
    for source_name in a["sources"]:
        source = safe_path(source_name, True)
        book = load_workbook(source, data_only=False)
        for sheet in book.worksheets:
            merged = output.create_sheet(unique_sheet_name(output, f"{source.stem}_{sheet.title}"))
            for row in sheet.iter_rows():
                for cell in row:
                    new_cell = merged[cell.coordinate]
                    new_cell.value = cell.value
                    if cell.has_style:
                        new_cell.font = copy.copy(cell.font)
                        new_cell.fill = copy.copy(cell.fill)
                        new_cell.border = copy.copy(cell.border)
                        new_cell.alignment = copy.copy(cell.alignment)
                        new_cell.protection = copy.copy(cell.protection)
                        new_cell.number_format = cell.number_format
            for region in sheet.merged_cells.ranges:
                merged.merge_cells(str(region))
        book.close()
    output.save(target)
    checked = load_workbook(target, data_only=False)
    sheet_names = checked.sheetnames
    checked.close()
    rendered = await validate_office(target, "openpyxl_reload")
    rendered["validation"]["sheets"] = sheet_names
    return {"ok": True, "artifact": artifact(target), **rendered}


async def analyze_xlsx(a: dict[str, Any]) -> dict[str, Any]:
    source = safe_path(a["path"], True)
    book = load_workbook(source, data_only=True, read_only=True)
    summary: dict[str, Any] = {}
    for sheet in book.worksheets:
        rows = list(sheet.iter_rows(values_only=True))
        headers = [str(value) if value is not None else f"column_{index}" for index, value in enumerate(rows[0], 1)] if rows else []
        numeric: dict[str, dict[str, float | int]] = {}
        for index, header in enumerate(headers):
            values = [row[index] for row in rows[1:] if index < len(row) and isinstance(row[index], (int, float)) and not isinstance(row[index], bool)]
            if values:
                numeric[header] = {"count": len(values), "sum": float(sum(values)), "average": float(sum(values) / len(values)), "min": float(min(values)), "max": float(max(values))}
        summary[sheet.title] = {"rows": max(len(rows) - 1, 0), "columns": len(headers), "headers": headers, "numeric": numeric}
    book.close()
    result: dict[str, Any] = {"ok": True, "analysis": summary}
    if a.get("output"):
        target = output_path(a["output"], ".xlsx")
        report = Workbook()
        overview = report.active
        overview.title = "分析汇总"
        overview.append(["Sheet", "数据行", "列数", "数值字段"])
        for name, data in summary.items():
            overview.append([name, data["rows"], data["columns"], ", ".join(data["numeric"])])
        overview.freeze_panes = "A2"
        overview.auto_filter.ref = overview.dimensions
        for cell in overview[1]:
            cell.font = Font(name="Noto Sans CJK SC", bold=True, color="FFFFFF")
            cell.fill = PatternFill("solid", fgColor="2F5597")
        report.save(target)
        checked = await validate_office(target, "openpyxl_reload")
        result.update({"artifact": artifact(target), **checked})
    return result


async def read_xlsx(a: dict[str, Any]) -> dict[str, Any]:
    path = safe_path(a["path"], True)
    book = load_workbook(path, data_only=a.get("data_only", False), read_only=True)
    result = {sheet.title: [list(row) for row in sheet.iter_rows(values_only=True)] for sheet in book.worksheets}
    book.close()
    return {"ok": True, "sheets": result}


def add_pptx_slide(presentation: Presentation, spec: dict[str, Any], title_slide: bool = False) -> None:
        layout_index = 0 if title_slide else min(1, len(presentation.slide_layouts) - 1)
        slide = presentation.slides.add_slide(presentation.slide_layouts[layout_index])
        slide.shapes.title.text = str(spec.get("title", ""))
        title_frame = slide.shapes.title.text_frame
        for paragraph in title_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = "Noto Sans CJK SC"
                run.font.size = PptPt(28)
                run.font.bold = True
        body = spec.get("bullets") or spec.get("content") or []
        if isinstance(body, str):
            body = [body]
        if len(slide.placeholders) > 1:
            frame = slide.placeholders[1].text_frame
            frame.clear()
            for line_index, text in enumerate(body):
                paragraph = frame.paragraphs[0] if line_index == 0 else frame.add_paragraph()
                paragraph.text = str(text)
                for run in paragraph.runs:
                    run.font.name = "Noto Sans CJK SC"
                    run.font.size = PptPt(20)
        if spec.get("image"):
            slide.shapes.add_picture(str(safe_path(spec["image"], True)), PptInches(7), PptInches(1.5), width=PptInches(5.5))
        table_data = spec.get("table") or []
        if table_data:
            rows = len(table_data)
            cols = max(len(row) for row in table_data)
            table = slide.shapes.add_table(rows, cols, PptInches(0.8), PptInches(2.2), PptInches(11.7), PptInches(3.8)).table
            for row_index, row in enumerate(table_data):
                for col_index, value in enumerate(row):
                    table.cell(row_index, col_index).text = str(value)
        chart_spec = spec.get("chart")
        if chart_spec:
            chart_data = ChartData()
            chart_data.categories = [str(value) for value in chart_spec.get("categories", [])]
            for series in chart_spec.get("series", []):
                chart_data.add_series(str(series.get("name", "数据")), [float(value) for value in series.get("values", [])])
            chart_types = {"bar": XL_CHART_TYPE.COLUMN_CLUSTERED, "line": XL_CHART_TYPE.LINE, "pie": XL_CHART_TYPE.PIE}
            chart_type = chart_types.get(str(chart_spec.get("type", "bar")).lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)
            slide.shapes.add_chart(chart_type, PptInches(1), PptInches(2), PptInches(11), PptInches(4.5), chart_data)
        for box in spec.get("boxes", []):
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                PptInches(float(box.get("x", 1))),
                PptInches(float(box.get("y", 2))),
                PptInches(float(box.get("width", 2.5))),
                PptInches(float(box.get("height", 1))),
            )
            shape.text = str(box.get("text", ""))
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "Noto Sans CJK SC"
                    run.font.size = PptPt(int(box.get("font_size", 16)))


async def create_pptx(a: dict[str, Any]) -> dict[str, Any]:
    path = output_path(a["path"], ".pptx")
    template = a.get("template")
    presentation = Presentation(str(safe_template("ppt", template))) if template else Presentation()
    presentation.slide_width = PptInches(13.333)
    presentation.slide_height = PptInches(7.5)
    for index, spec in enumerate(a["slides"]):
        add_pptx_slide(presentation, spec, title_slide=index == 0 and not template)
    presentation.save(path)
    checked = Presentation(path)
    count = len(checked.slides)
    rendered = await validate_office(path, "python_pptx_reload")
    rendered["validation"]["slides"] = count
    return {"ok": True, "artifact": artifact(path), **rendered}


async def edit_pptx(a: dict[str, Any]) -> dict[str, Any]:
    source = safe_path(a.get("source") or a["path"], True)
    target = output_path(a.get("output") or a.get("path") or source.name, ".pptx")
    presentation = Presentation(source)
    replacements = {str(key): str(value) for key, value in a.get("replace", {}).items()}
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                original = paragraph.text
                updated = original
                for old, new in replacements.items():
                    updated = updated.replace(old, new)
                if updated != original:
                    paragraph.text = updated
    for spec in a.get("append_slides", []):
        add_pptx_slide(presentation, spec)
    presentation.save(target)
    checked = Presentation(target)
    count = len(checked.slides)
    rendered = await validate_office(target, "python_pptx_reload")
    rendered["validation"]["slides"] = count
    return {"ok": True, "artifact": artifact(target), **rendered}


async def read_pptx(a: dict[str, Any]) -> dict[str, Any]:
    presentation = Presentation(safe_path(a["path"], True))
    slides = []
    for slide in presentation.slides:
        slides.append([shape.text for shape in slide.shapes if hasattr(shape, "text_frame") and shape.has_text_frame])
    return {"ok": True, "slides": slides}


async def office_to_pdf(a: dict[str, Any]) -> dict[str, Any]:
    source = safe_path(a["path"], True)
    pdf = await libreoffice_convert(source)
    with fitz.open(pdf) as opened:
        pages = opened.page_count
    return {"ok": True, "artifact": artifact(pdf), "validation": {"pymupdf_open": "PASS", "pages": pages}}


async def read_pdf(a: dict[str, Any]) -> dict[str, Any]:
    path = safe_path(a["path"], True)
    with fitz.open(path) as document:
        pages = [page.get_text() for page in document]
    return {"ok": True, "pages": pages, "page_count": len(pages)}


async def pdf_to_images(a: dict[str, Any]) -> dict[str, Any]:
    source = safe_path(a["path"], True)
    dpi = max(72, min(int(a.get("dpi", 150)), 600))
    output_dir = safe_path(a.get("output_dir", f"outputs/{source.stem}-pages"))
    output_dir.mkdir(parents=True, exist_ok=True)
    artifacts = []
    with fitz.open(source) as document:
        for index, page in enumerate(document, 1):
            target = output_dir / f"page-{index:04d}.png"
            page.get_pixmap(dpi=dpi, alpha=False).save(target)
            artifacts.append(artifact(target))
    return {"ok": True, "artifacts": artifacts, "page_count": len(artifacts), "dpi": dpi}


async def merge_pdf(a: dict[str, Any]) -> dict[str, Any]:
    target = output_path(a["path"], ".pdf")
    writer = PdfWriter()
    for name in a["sources"]:
        for page in PdfReader(safe_path(name, True)).pages:
            writer.add_page(page)
    with target.open("wb") as handle:
        writer.write(handle)
    return {"ok": True, "artifact": artifact(target)}


async def split_pdf(a: dict[str, Any]) -> dict[str, Any]:
    source = safe_path(a["path"], True)
    reader = PdfReader(source)
    artifacts = []
    for index, page in enumerate(reader.pages, 1):
        target = output_path(f"{source.stem}-page-{index}.pdf", ".pdf")
        writer = PdfWriter()
        writer.add_page(page)
        with target.open("wb") as handle:
            writer.write(handle)
        artifacts.append(artifact(target))
    return {"ok": True, "artifacts": artifacts}


HANDLERS = {
    "create_docx": create_docx,
    "edit_docx": edit_docx,
    "read_docx": read_docx,
    "docx_to_pdf": office_to_pdf,
    "create_xlsx": create_xlsx,
    "edit_xlsx": edit_xlsx,
    "merge_xlsx": merge_xlsx,
    "analyze_xlsx": analyze_xlsx,
    "read_xlsx": read_xlsx,
    "create_pptx": create_pptx,
    "edit_pptx": edit_pptx,
    "read_pptx": read_pptx,
    "pptx_to_pdf": office_to_pdf,
    "office_to_pdf": office_to_pdf,
    "read_pdf": read_pdf,
    "extract_pdf_text": read_pdf,
    "pdf_to_images": pdf_to_images,
    "merge_pdf": merge_pdf,
    "split_pdf": split_pdf,
}


@app.get("/health")
async def health() -> dict[str, Any]:
    return {"status": "ok", "libreoffice": shutil.which("libreoffice") is not None, "fonts": ["Noto Sans CJK SC", "Noto Serif CJK SC"]}


@app.post("/tools/execute")
async def execute(request: ToolRequest, authorization: str | None = Header(default=None)) -> dict[str, Any]:
    authorize(authorization)
    try:
        if request.name not in HANDLERS:
            raise KeyError(f"unknown office tool: {request.name}")
        return await HANDLERS[request.name](request.arguments)
    except Exception as exc:
        return {"ok": False, "error": f"{type(exc).__name__}: {exc}"}
